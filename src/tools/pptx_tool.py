"""PPTXTool - Wrapper for python-pptx library.

This module provides a clean API for creating PowerPoint presentations
programmatically using the python-pptx library.
"""

import logging
from pathlib import Path
from typing import TYPE_CHECKING

from pptx import Presentation
from pptx.util import Inches

if TYPE_CHECKING:
    from pptx.slide import Slide

logger = logging.getLogger(__name__)


class PPTXToolError(Exception):
    """Base exception for PPTXTool errors."""

    pass


class PPTXTool:
    """Tool for creating and manipulating PowerPoint presentations.

    This tool provides a clean API for creating PowerPoint presentations
    programmatically, supporting slides with different layouts, text,
    images, charts, and tables.

    Args:
        template_path: Optional path to a template PPTX file. If None,
            creates a blank presentation.

    Example:
        >>> tool = PPTXTool()
        >>> slide = tool.add_slide("title_and_content")
        >>> tool.add_title(slide, "My Title")
        >>> tool.add_bullets(slide, ["Point 1", "Point 2"])
        >>> tool.save("output.pptx")
    """

    def __init__(self, template_path: str | None = None) -> None:
        """Initialize PPTXTool with new or template presentation."""
        try:
            if template_path:
                template_path_obj = Path(template_path)
                if not template_path_obj.exists():
                    raise PPTXToolError(f"Template file not found: {template_path}")
                logger.debug(f"Loading template from: {template_path}")
                self.presentation = Presentation(str(template_path))
            else:
                logger.debug("Creating new blank presentation")
                self.presentation = Presentation()
        except PPTXToolError:
            raise
        except Exception as e:
            logger.error(f"Failed to initialize presentation: {e}")
            raise PPTXToolError(f"Failed to initialize presentation: {e}") from e

    def add_slide(self, layout: str) -> "Slide":
        """Add a new slide with specified layout.

        Args:
            layout: Layout type for the slide. Options:
                - "title": Title slide layout
                - "title_and_content": Title with bullet points
                - "section_header": Section header layout
                - "blank": Blank slide

        Returns:
            The created Slide object.

        Raises:
            PPTXToolError: If layout is not recognized.
        """
        layout_map = {
            "title": 0,  # Title Slide
            "title_and_content": 1,  # Title and Content
            "section_header": 2,  # Section Header
            "blank": 6,  # Blank
        }

        if layout not in layout_map:
            raise PPTXToolError(
                f"Unknown layout: {layout}. Must be one of {list(layout_map.keys())}"
            )

        try:
            layout_idx = layout_map[layout]
            prs_layout = self.presentation.slide_layouts[layout_idx]
            slide = self.presentation.slides.add_slide(prs_layout)
            logger.debug(f"Added slide with layout: {layout}")
            return slide  # type: ignore[no-any-return]
        except Exception as e:
            logger.error(f"Failed to add slide with layout {layout}: {e}")
            raise PPTXToolError(f"Failed to add slide: {e}") from e

    def add_title(self, slide: "Slide", text: str) -> None:
        """Add title to slide.

        Args:
            slide: The Slide object to add title to.
            text: The title text.

        Raises:
            PPTXToolError: If slide has no title placeholder.
        """
        try:
            if slide.shapes.title is None:
                raise PPTXToolError("Slide does not have a title placeholder")
            slide.shapes.title.text = text
            logger.debug(f"Added title: {text[:50]}...")
        except PPTXToolError:
            raise
        except Exception as e:
            logger.error(f"Failed to add title: {e}")
            raise PPTXToolError(f"Failed to add title: {e}") from e

    def add_bullets(  # noqa: C901 - Complex function for handling various bullet cases
        self,
        slide: "Slide",
        items: list[str],
        level: int = 0,
    ) -> None:
        """Add bullet points to slide.

        Args:
            slide: The Slide object to add bullets to.
            items: List of text items to add as bullets.
            level: Indentation level (0 for top level, 1, 2, etc. for nested).

        Raises:
            PPTXToolError: If slide has no body placeholder.
        """
        try:
            # Find body shape
            body_shape = None
            for shape in slide.shapes:
                if shape.has_text_frame and hasattr(shape, "placeholder_format"):
                    try:
                        if shape.placeholder_format.idx == 1:  # Body placeholder
                            body_shape = shape
                            break
                    except Exception:
                        continue

            if body_shape is None:
                # Try to find any shape with text frame
                for shape in slide.shapes:
                    if shape.has_text_frame and shape != slide.shapes.title:
                        body_shape = shape
                        break

            if body_shape is None:
                raise PPTXToolError("Slide does not have a body placeholder")

            text_frame = body_shape.text_frame  # type: ignore[attr-defined]
            text_frame.clear()

            for item in items:
                p = text_frame.add_paragraph()
                p.text = item
                p.level = level
                logger.debug(f"Added bullet: {item[:50]}... (level: {level})")
        except PPTXToolError:
            raise
        except Exception as e:
            logger.error(f"Failed to add bullets: {e}")
            raise PPTXToolError(f"Failed to add bullets: {e}") from e

    def add_image(
        self,
        slide: "Slide",
        image_path: str,
        left: float,
        top: float,
        width: float | None = None,
        height: float | None = None,
    ) -> None:
        """Add image to slide.

        Args:
            slide: The Slide object to add image to.
            image_path: Path to the image file.
            left: Left position in inches.
            top: Top position in inches.
            width: Optional width in inches. If not provided, uses image width.
            height: Optional height in inches. If not provided, uses image height.

        Raises:
            PPTXToolError: If image file not found or invalid.
        """
        try:
            image_path_obj = Path(image_path)
            if not image_path_obj.exists():
                raise PPTXToolError(f"Image file not found: {image_path}")

            # Convert to inches
            left_inches = Inches(left)
            top_inches = Inches(top)
            width_inches = Inches(width) if width else None
            height_inches = Inches(height) if height else None

            # Add picture to slide
            if width_inches and height_inches:
                slide.shapes.add_picture(
                    str(image_path_obj),
                    left_inches,
                    top_inches,
                    width=width_inches,
                    height=height_inches,
                )
            elif width_inches:
                slide.shapes.add_picture(
                    str(image_path_obj),
                    left_inches,
                    top_inches,
                    width=width_inches,
                )
            elif height_inches:
                slide.shapes.add_picture(
                    str(image_path_obj),
                    left_inches,
                    top_inches,
                    height=height_inches,
                )
            else:
                slide.shapes.add_picture(
                    str(image_path_obj),
                    left_inches,
                    top_inches,
                )

            logger.debug(
                f"Added image: {image_path} (position: {left}, {top}; size: {width}x{height})"
            )
        except PPTXToolError:
            raise
        except Exception as e:
            logger.error(f"Failed to add image: {e}")
            raise PPTXToolError(f"Failed to add image: {e}") from e

    def add_table(
        self,
        slide: "Slide",
        rows: int,
        cols: int,
        data: list[list[str]],
        left: float,
        top: float,
        width: float,
        height: float,
    ) -> None:
        """Add table to slide.

        Args:
            slide: The Slide object to add table to.
            rows: Number of rows in table.
            cols: Number of columns in table.
            data: 2D list of cell data. Should have rows x cols elements.
            left: Left position in inches.
            top: Top position in inches.
            width: Table width in inches.
            height: Table height in inches.

        Raises:
            PPTXToolError: If data dimensions don't match rows/cols.
        """
        try:
            # Validate data dimensions
            if len(data) != rows:
                raise PPTXToolError(f"Data has {len(data)} rows but {rows} expected")
            for i, row in enumerate(data):
                if len(row) != cols:
                    raise PPTXToolError(f"Row {i} has {len(row)} columns but {cols} expected")

            # Convert to inches
            left_inches = Inches(left)
            top_inches = Inches(top)
            width_inches = Inches(width)
            height_inches = Inches(height)

            # Add table shape
            table_shape = slide.shapes.add_table(
                rows,
                cols,
                left_inches,
                top_inches,
                width_inches,
                height_inches,
            ).table

            # Fill in table data
            for row_idx in range(rows):
                for col_idx in range(cols):
                    cell = table_shape.cell(row_idx, col_idx)
                    cell.text = data[row_idx][col_idx]

            logger.debug(
                f"Added table: {rows}x{cols} (position: {left}, {top}; size: {width}x{height})"
            )
        except PPTXToolError:
            raise
        except Exception as e:
            logger.error(f"Failed to add table: {e}")
            raise PPTXToolError(f"Failed to add table: {e}") from e

    def save(self, output_path: str) -> None:
        """Save presentation to file.

        Args:
            output_path: Path where to save the presentation file.

        Raises:
            PPTXToolError: If save operation fails.
        """
        try:
            output_path_obj = Path(output_path)

            # Ensure parent directory exists
            output_path_obj.parent.mkdir(parents=True, exist_ok=True)

            # Save presentation
            self.presentation.save(str(output_path_obj))
            logger.info(f"Presentation saved to: {output_path}")
        except Exception as e:
            logger.error(f"Failed to save presentation: {e}")
            raise PPTXToolError(f"Failed to save presentation: {e}") from e
