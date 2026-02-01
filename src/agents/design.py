"""Design agent for assembling PowerPoint slides."""

from __future__ import annotations

from typing import TYPE_CHECKING

from src.models.schemas import SlideContent
from src.tools.pptx_tool import PPTXTool

if TYPE_CHECKING:
    from pptx.presentation import Presentation
    from pptx.slide import Slide


class DesignAgentError(RuntimeError):
    """Raised when the design agent fails to build slides."""


class DesignAgent:
    """Creates PowerPoint slides from generated content."""

    def __init__(self, template_path: str | None = None) -> None:
        self.pptx_tool = PPTXTool(template_path=template_path)

    def create_slide(
        self,
        content: SlideContent,
        layout: str | None = None,
    ) -> Slide:
        """Create a single slide from SlideContent."""
        resolved_layout = layout
        if resolved_layout is None:
            resolved_layout = "title" if not content.content else "title_and_content"

        try:
            slide = self.pptx_tool.add_slide(resolved_layout)
            self.pptx_tool.add_title(slide, content.title)

            if content.content:
                self.pptx_tool.add_bullets(slide, content.content)

            if content.notes and slide.notes_slide and slide.notes_slide.notes_text_frame:
                slide.notes_slide.notes_text_frame.text = content.notes

            return slide
        except Exception as exc:  # noqa: BLE001 - wrap any creation failure
            raise DesignAgentError("Failed to create slide") from exc

    def create_presentation(
        self,
        slides: list[SlideContent],
        output_path: str | None = None,
    ) -> Presentation:
        """Create a presentation from slide contents."""
        try:
            for content in slides:
                self.create_slide(content)

            if output_path:
                self.pptx_tool.save(output_path)

            return self.pptx_tool.presentation
        except Exception as exc:  # noqa: BLE001 - wrap any creation failure
            raise DesignAgentError("Failed to create presentation") from exc
